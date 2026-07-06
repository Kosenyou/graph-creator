import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # 16:9 ワイドスクリーンに設定
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # カラーパレットの定義 (プレミアム感のあるネイビーとブルー)
    COLOR_PRIMARY = RGBColor(26, 54, 93)      # 濃紺 (#1A365D)
    COLOR_SECONDARY = RGBColor(43, 108, 176)   # ブルー (#2B6CB0)
    COLOR_ACCENT = RGBColor(49, 151, 149)     # ティール (#319795)
    COLOR_BG_LIGHT = RGBColor(247, 250, 252)  # 薄いグレー (#F7FAFC)
    COLOR_TEXT_DARK = RGBColor(45, 55, 72)     # ダークグレー (#2D3748)
    COLOR_TEXT_MUTED = RGBColor(113, 128, 150) # ミューテッドグレー (#718096)
    COLOR_WHITE = RGBColor(255, 255, 255)
    
    # フォント設定ヘルパー
    def set_font(run, name="Yu Gothic", size=Pt(14), bold=False, italic=False, color=COLOR_TEXT_DARK):
        run.font.name = name
        run.font.size = size
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color

    # 背景を設定するヘルパー
    def add_background(slide, color=COLOR_BG_LIGHT):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # スライドタイトルとヘッダーラインを追加するヘルパー
    def add_slide_header(slide, title_text):
        add_background(slide)
        
        # タイトルテキストボックス
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        run = p.runs[0]
        set_font(run, name="Meiryo", size=Pt(28), bold=True, color=COLOR_PRIMARY)
        
        # 装飾用の下線（極細の長方形）
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_SECONDARY
        line.line.fill.background() # 枠線なし

    # ----------------------------------------------------
    # スライド 1: タイトルスライド (ダーク背景でプレミアム感を演出)
    # ----------------------------------------------------
    blank_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_layout)
    add_background(slide1, COLOR_PRIMARY)
    
    # 装飾用のグラデーション風アクセントシェイプ (左側に配置)
    accent_bar = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.4), Inches(7.5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.fill.background()
    
    # タイトルとサブタイトルの配置
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "論文用グラフ作成"
    p1.alignment = PP_ALIGN.LEFT
    run1 = p1.runs[0]
    set_font(run1, name="Meiryo", size=Pt(48), bold=True, color=COLOR_WHITE)
    
    p2 = tf.add_paragraph()
    p2.text = "CSV & Excel 論文用グラフ作成・エクスポートツール"
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(15)
    run2 = p2.runs[0]
    set_font(run2, name="Meiryo", size=Pt(20), bold=False, color=COLOR_WHITE)

    p3 = tf.add_paragraph()
    p3.text = "機能紹介 ＆ 取扱説明書"
    p3.alignment = PP_ALIGN.LEFT
    p3.space_before = Pt(10)
    run3 = p3.runs[0]
    set_font(run3, name="Meiryo", size=Pt(20), bold=True, color=COLOR_ACCENT)

    p4 = tf.add_paragraph()
    p4.text = "Free Edition (多言語対応: 日本語 / English)"
    p4.alignment = PP_ALIGN.LEFT
    p4.space_before = Pt(40)
    run4 = p4.runs[0]
    set_font(run4, name="Yu Gothic", size=Pt(12), color=COLOR_WHITE)

    # ----------------------------------------------------
    # スライド 2: 製品概要 (Overview)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide2, "1. 製品概要 (Overview)")
    
    # 左側の紹介テキスト
    desc_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(7.2), Inches(4.8))
    tf = desc_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "「論文用グラフ作成」は、手元のCSVやExcelファイルをアップロードするだけで、美しく直感的なグラフを瞬時に作成・保存できるWebアプリケーションです。"
    set_font(p.runs[0], size=Pt(18), bold=True, color=COLOR_TEXT_DARK)
    
    points = [
        "インストール不要: Webブラウザ上で完結するため、OSを問わず即座に利用可能。",
        "直感的なUI: プログラミングや複雑なExcel操作を必要とせず、マウス操作だけでグラフを作成。",
        "多様な保存形式: プレゼン用の画像保存から、データ付きExcelファイルの出力まで幅広く対応。",
        "多言語切り替え: 日本語と英語の表示切替に対応し、グローバルな環境でも使用可能。"
    ]
    for pt in points:
        p = tf.add_paragraph()
        p.text = "• " + pt
        p.space_before = Pt(14)
        p.space_after = Pt(2)
        set_font(p.runs[0], size=Pt(14), color=COLOR_TEXT_DARK)

    # 右側の目を引くハイライトボックス
    box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(1.8), Inches(4.0), Inches(4.8)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_PRIMARY
    box.line.fill.background()
    
    box_tf = box.text_frame
    box_tf.word_wrap = True
    box_tf.margin_left = Inches(0.4)
    box_tf.margin_right = Inches(0.4)
    box_tf.margin_top = Inches(0.5)
    
    bp1 = box_tf.paragraphs[0]
    bp1.text = "主な特徴と提供価値"
    set_font(bp1.runs[0], size=Pt(20), bold=True, color=COLOR_WHITE)
    
    features = [
        "🚀 瞬時のビジュアライズ",
        "📂 Excelの複数表の自動認識",
        "📊 高度な2軸折れ線・散布図対応",
        "💾 グラフ画像埋込済Excelの出力"
    ]
    for feat in features:
        bp = box_tf.add_paragraph()
        bp.text = feat
        bp.space_before = Pt(20)
        set_font(bp.runs[0], size=Pt(15), bold=True, color=COLOR_WHITE)

    # ----------------------------------------------------
    # スライド 3: ファイル読み込み機能
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide3, "2. 高性能なファイル読み込み機能")
    
    # 左右に並べる2つのカードボックス
    # カード1: CSV読み込み
    card1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_WHITE
    card1.line.color.rgb = COLOR_SECONDARY
    card1.line.width = Pt(1.5)
    
    c1_tf = card1.text_frame
    c1_tf.word_wrap = True
    c1_tf.margin_left = Inches(0.3)
    c1_tf.margin_right = Inches(0.3)
    c1_tf.margin_top = Inches(0.3)
    
    p = c1_tf.paragraphs[0]
    p.text = "CSVファイルの読み込み"
    set_font(p.runs[0], size=Pt(20), bold=True, color=COLOR_PRIMARY)
    
    csv_features = [
        "シンプルな区切り文字フォーマットを高速パース。",
        "1行目を列ヘッダー（列名）として自動判定し、軸の選択肢に反映。",
        "カンマ区切りの数値データを素早くプロット。"
    ]
    for cf in csv_features:
        p = c1_tf.add_paragraph()
        p.text = "✔ " + cf
        p.space_before = Pt(14)
        set_font(p.runs[0], size=Pt(14), color=COLOR_TEXT_DARK)

    # カード2: Excel読み込み (プレミアム機能)
    card2 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    card2.fill.solid()
    card2.fill.fore_color.rgb = COLOR_WHITE
    card2.line.color.rgb = COLOR_ACCENT
    card2.line.width = Pt(1.5)
    
    c2_tf = card2.text_frame
    c2_tf.word_wrap = True
    c2_tf.margin_left = Inches(0.3)
    c2_tf.margin_right = Inches(0.3)
    c2_tf.margin_top = Inches(0.3)
    
    p = c2_tf.paragraphs[0]
    p.text = "Excelファイル (.xlsx / .xlsm) 対応"
    set_font(p.runs[0], size=Pt(20), bold=True, color=COLOR_ACCENT)
    
    excel_features = [
        "複数表の自動検出: 1シート内に点在する複数の表を数値データのまとまりから自動検知して個別に読み込み可能。",
        "複雑な見出しの結合解除: セルの結合や複数段の見出しを自動解析し、[見出し1 / 見出し2] のように列名を自動生成。",
        "不要な計算行の除外: データの末尾にある「平均」や「標準偏差」といった統計計算行を自動判定し、グラフのプロット範囲から除外。"
    ]
    for ef in excel_features:
        p = c2_tf.add_paragraph()
        p.text = "✔ " + ef
        p.space_before = Pt(12)
        set_font(p.runs[0], size=Pt(13), color=COLOR_TEXT_DARK)

    # ----------------------------------------------------
    # スライド 4: グラフ表現とカスタマイズ
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide4, "3. 多彩なグラフ表現と柔軟なカスタマイズ")
    
    # 3つの列でグラフ種類を紹介
    g_types = [
        ("折れ線グラフ (Line)", "時間の経過に伴う推移や連続的なデータのトレンド変化を視覚化するのに最適です。", COLOR_PRIMARY),
        ("散布図 (Scatter)", "2つの数値軸（X軸・Y軸）の相関関係を分析します。両軸に数値データが必要です。", COLOR_SECONDARY),
        ("棒グラフ (Bar)", "項目間の数量比較や、不連続なカテゴリデータの可視化に威力を発揮します。", COLOR_ACCENT)
    ]
    
    for idx, (title, desc, color) in enumerate(g_types):
        left_pos = Inches(0.8 + idx * 4.0)
        col_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8), Inches(3.7), Inches(2.2))
        col_box.fill.solid()
        col_box.fill.fore_color.rgb = color
        col_box.line.fill.background()
        
        c_tf = col_box.text_frame
        c_tf.word_wrap = True
        c_tf.margin_left = Inches(0.2)
        c_tf.margin_right = Inches(0.2)
        c_tf.margin_top = Inches(0.2)
        
        p = c_tf.paragraphs[0]
        p.text = title
        set_font(p.runs[0], size=Pt(16), bold=True, color=COLOR_WHITE)
        
        p = c_tf.add_paragraph()
        p.text = desc
        p.space_before = Pt(8)
        set_font(p.runs[0], size=Pt(12), color=COLOR_WHITE)

    # 下部領域にカスタマイズ機能をまとめる
    cust_box = slide4.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.5))
    ctf = cust_box.text_frame
    ctf.word_wrap = True
    
    p = ctf.paragraphs[0]
    p.text = "グラフィックとレイアウトの細やかな制御項目"
    set_font(p.runs[0], size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    c_details = [
        "複数系列 (Multi-Series) 対応: 「+ 系列を追加」ボタンから、1つのグラフ上に異なるY列のデータを何個でも重ね描きできます。",
        "左右2軸対応: 単位や桁数が大きく異なる指標を比較するため、「左Y軸」と「右Y軸」を個別に設定可能です。",
        "表示オプション: 凡例の位置変更（右上 / 左上 / 右下 / 左下）、および「X軸の反転」機能を搭載しています。"
    ]
    for cd in c_details:
        p = ctf.add_paragraph()
        p.text = "• " + cd
        p.space_before = Pt(8)
        set_font(p.runs[0], size=Pt(13), color=COLOR_TEXT_DARK)

    # ----------------------------------------------------
    # スライド 5: 保存・エクスポート機能
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide5, "4. 強力なエクスポート機能")
    
    exports = [
        ("📸 PNG保存", "高解像度の画像形式としてグラフをエクスポート。資料や報告書への貼り付けに最適です。"),
        ("📄 PDF保存", "描画されたグラフをPDF形式に直接変換して保存します。印刷や共有に適したレイアウトで出力されます。"),
        ("📊 XLSX保存", "読み込んだ元の表データと、作成したグラフ画像をセットにして1つのExcelファイルとして保存します。"),
        ("⚡ XLSM保存", "マクロ有効ブック形式 (.xlsm) として出力します。マクロが含まれる業務テンプレート等に組み込むベースとして使用可能です（本アプリからはマクロ本体は含めず、安全に保存されます）。")
    ]
    
    for idx, (title, desc) in enumerate(exports):
        row_pos = Inches(1.8 + idx * 1.3)
        # アイコン用ボックス
        t_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), row_pos, Inches(2.5), Inches(1.0))
        t_box.fill.solid()
        t_box.fill.fore_color.rgb = COLOR_PRIMARY if idx < 2 else COLOR_ACCENT
        t_box.line.fill.background()
        
        ttf = t_box.text_frame
        ttf.word_wrap = True
        ttf.margin_top = Inches(0.3)
        p = ttf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        set_font(p.runs[0], size=Pt(16), bold=True, color=COLOR_WHITE)
        
        # 説明用ボックス
        d_box = slide5.shapes.add_textbox(Inches(3.5), row_pos - Inches(0.1), Inches(9.0), Inches(1.0))
        dtf = d_box.text_frame
        dtf.word_wrap = True
        p = dtf.paragraphs[0]
        p.text = desc
        set_font(p.runs[0], size=Pt(14), color=COLOR_TEXT_DARK)

    # ----------------------------------------------------
    # スライド 6: かんたん操作手順 (How to Use)
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide6, "5. 操作手順 (How to Use)")
    
    steps = [
        ("Step 1", "ファイルの読込", "「CSV/Excel読み込み」ボタンから、お手元のデータファイルを選択します。"),
        ("Step 2", "データの指定", "（Excelの場合）対象シートと、自動検出された表候補を選択します。"),
        ("Step 3", "グラフの設定", "グラフ種類、X軸の列、Y軸の列（複数可）、タイトル、軸ラベルを設定します。"),
        ("Step 4", "作成と保存", "「グラフ作成」ボタンを押し、表示を確認したらお好みの形式で保存します。")
    ]
    
    for idx, (num, step_name, step_desc) in enumerate(steps):
        left_pos = Inches(0.8 + idx * 2.95)
        # ステップ四角形
        step_shape = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8), Inches(2.8), Inches(4.8))
        step_shape.fill.solid()
        step_shape.fill.fore_color.rgb = COLOR_WHITE
        step_shape.line.color.rgb = COLOR_SECONDARY
        step_shape.line.width = Pt(1.5)
        
        stf = step_shape.text_frame
        stf.word_wrap = True
        stf.margin_left = Inches(0.2)
        stf.margin_right = Inches(0.2)
        stf.margin_top = Inches(0.3)
        
        # 番号
        p = stf.paragraphs[0]
        p.text = num
        p.alignment = PP_ALIGN.CENTER
        set_font(p.runs[0], size=Pt(24), bold=True, color=COLOR_ACCENT)
        
        # 名前
        p = stf.add_paragraph()
        p.text = step_name
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(10)
        set_font(p.runs[0], size=Pt(16), bold=True, color=COLOR_PRIMARY)
        
        # 装飾横線
        p = stf.add_paragraph()
        p.text = "━━━━━━"
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(5)
        set_font(p.runs[0], size=Pt(10), color=COLOR_TEXT_MUTED)
        
        # 説明文
        p = stf.add_paragraph()
        p.text = step_desc
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(14)
        set_font(p.runs[0], size=Pt(12), color=COLOR_TEXT_DARK)

    # ----------------------------------------------------
    # スライド 7: 結び (End)
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_background(slide7, COLOR_PRIMARY)
    
    end_box = slide7.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(11.0), Inches(3.0))
    etf = end_box.text_frame
    etf.word_wrap = True
    
    p = etf.paragraphs[0]
    p.text = "論文用グラフ作成"
    p.alignment = PP_ALIGN.CENTER
    set_font(p.runs[0], size=Pt(44), bold=True, color=COLOR_WHITE)
    
    p = etf.add_paragraph()
    p.text = "誰でもかんたんに、すばやく、美しいデータをレポートへ。"
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(15)
    set_font(p.runs[0], size=Pt(18), color=COLOR_WHITE)

    p = etf.add_paragraph()
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(40)
    set_font(p.runs[0], size=Pt(28), bold=True, color=COLOR_ACCENT)

    # 保存
    output_path = "c:\\Users\\forSc\\OneDrive\\ドキュメント\\グラフ作成_NoAI\\論文用グラフ作成_機能紹介.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
